import sys, os
sys.path.append(os.path.join(os.path.dirname(__file__), ".."))
from PIL import Image
import collections.abc
from typing import Optional
from fastapi import FastAPI
from pptx.util import Inches
from pptx import Presentation
from api_key import *
from bing_image_downloader import downloader
from bing_image_urls import bing_image_urls

from src.utils import *
from src.text_gen import *

import requests

import random
import io
import urllib
from fastapi.responses import FileResponse



DATA_FOLDER = r"data"
FONT_FOLDER = r"fonts"
IMAGE_FOLDER = os.path.join(r"images")
TEMPLATE_PPTX = os.path.join(DATA_FOLDER, "template.pptx")
CHOSEN_FONT = os.path.join(FONT_FOLDER, "Calibri Regular.ttf")

from fastapi.middleware.cors import CORSMiddleware

app = FastAPI()

# middlewares
app.add_middleware(
    CORSMiddleware, # https://fastapi.tiangolo.com/tutorial/cors/
    allow_origins=['*'], # wildcard to allow all, more here - https://developer.mozilla.org/en-US/docs/Web/HTTP/Headers/Access-Control-Allow-Origin
    allow_credentials=True, # https://developer.mozilla.org/en-US/docs/Web/HTTP/Headers/Access-Control-Allow-Credentials
    allow_methods=['*'], # https://developer.mozilla.org/en-US/docs/Web/HTTP/Headers/Access-Control-Allow-Methods
    allow_headers=['*'], # https://developer.mozilla.org/en-US/docs/Web/HTTP/Headers/Access-Control-Allow-Headers
)

# import nest_asyncio
# from pyngrok import ngrok
# import uvicorn

# ngrok.set_auth_token(NGROK_TOKEN)

# # specify a port
# port = 8000
# ngrok_tunnel = ngrok.connect(port)

# # where we can visit our fastAPI app
# print('Public URL:', ngrok_tunnel.public_url)


# nest_asyncio.apply()




@app.post("/generate/")
async def generate(topic: str, mode: int = 0, n_slides: Optional[int] = 10, n_words_per_slide: Optional[int] = 70):
    """ """
    if mode == 0:
        text_query = create_query_from_text(topic, type_of_text="topic", n_slides=n_slides, n_words_per_slide=n_words_per_slide)
    else:
        text_query = create_query_from_text(topic, type_of_text="doc", n_slides=n_slides, n_words_per_slide=n_words_per_slide)

    response = query_from_API(query=text_query, token=POE_API_KEY)
    content_json = create_content_json(response)

    if content_json is None:
        print("Cannot extract json from text. Cannot create a presentation. Stopped.")
        exit()

    try:
        default_16_9_slide_size = (Inches(5.625), Inches(10))
        prs = Presentation(TEMPLATE_PPTX)
        if not (
            prs.slide_height == default_16_9_slide_size[0] * 914400 and prs.slide_width == default_16_9_slide_size[1] * 914400
        ):
            print(f"Use template from {TEMPLATE_PPTX}")
        else:
            prs = Presentation()
            prs.slide_height, prs.slide_width = default_16_9_slide_size
            print("Template is not of 16:9 ratio, creating from blank template.")

    except:
        prs = Presentation()
        prs.slide_height, prs.slide_width = default_16_9_slide_size
        print(f"Cannot use template from {TEMPLATE_PPTX}. Creating a blank file.")

    layout_id = get_layout_id(prs)

    for i in range(len(prs.slides) - 1, -1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]

    # get_title_query = create_query_get_title(content_json)
    # title = query_from_API(query=get_title_query, token=POE_API_KEY)
    title = topic
    slide_title = prs.slides.add_slide(prs.slide_layouts[0])
    title_box = slide_title.shapes.title
    title_box.text = title.capitalize()
    for i, place_holder in enumerate(slide_title.placeholders):
        if i < 1:
            continue
        sp = place_holder.element
        sp.getparent().remove(sp)

    key = list(content_json.keys())[0]

    if not os.path.isdir(IMAGE_FOLDER):
        os.makedirs(IMAGE_FOLDER)

    img_slot = {"left": Inches(6), "top": Inches(1.6), "width": Inches(3.6), "height": Inches(3.6)}
    img_slot_ratio = img_slot["width"] / img_slot["height"]

    for item in content_json[key]:
        header, content = process_header(item["header"]), item["content"]
        image_query = f"{header}_{topic}".replace(" ", "_")
        image_pil = None
        if "Introduction" not in image_query:
            try:
                # downloader.download(
                #     image_query,
                #     limit=1,
                #     output_dir=IMAGE_FOLDER,
                #     force_replace=False,
                #     filter="photo",
                #     timeout=10,
                #     verbose=False,
                # )

                # image_folder_path = os.path.join(IMAGE_FOLDER, image_query)
                # image_path = os.path.join(image_folder_path, os.listdir(image_folder_path)[0])
                # image = Image.open(image_path)
                image_urls = bing_image_urls(image_query, limit=100)
                random.shuffle(image_urls)
                for url in image_urls:
                    try:
                        image_content = urllib.request.urlopen(url, timeout=3)
                        image_pil = Image.open(image_content)
                        image_out = io.BytesIO()
                        image_pil.save(image_out, format="PNG")
                        break
                    except:
                        print("Cannot download image from url: {}".format(url))
                        continue

                print('Downloaded image {} for slide "{}"'.format(url, header))
            except Exception as e:
                print(e)

        slide_layout = prs.slide_layouts[layout_id]
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = header

        for i, place_holder in enumerate(slide.placeholders):
            if i < 1:
                continue
            sp = place_holder.element
            sp.getparent().remove(sp)

        if image_pil:
            w, h = image_pil.size
            try:
                if w / h > img_slot_ratio:  # image longer than slot -> resize image to fit slot_width
                    picture = slide.shapes.add_picture(image_out, img_slot["left"], img_slot["top"], width=img_slot["width"])
                else:
                    picture = slide.shapes.add_picture(image_out, img_slot["left"], img_slot["top"], height=img_slot["height"])
            except Exception as e:
                print("Cannot add picture. ", e)

        left, top, width, height = Inches(1), Inches(1.5), Inches(5), Inches(5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = content
        text_frame.fit_text(font_file=CHOSEN_FONT, max_size=22)
        text_frame.word_wrap = True

    output_pptx_path = os.path.join("data", topic.replace(" ", "_") + ".pptx")
    output_pptx_path = change_name_if_duplicated(output_pptx_path)
    prs.save(output_pptx_path)

    output_folder = os.path.splitext(output_pptx_path)[-2]
    os.makedirs(output_folder)
    # convert_pptx_to_svg(pptx_file=output_pptx_path, output_folder=output_folder)
    print(f"Presentation saved to {output_folder}")

    
    return FileResponse(output_pptx_path, filename=output_pptx_path.split("/")[-1], media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")

# # finally run the app
# uvicorn.run(app, port=port)