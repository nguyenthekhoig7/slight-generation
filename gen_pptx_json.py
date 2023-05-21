import json
import collections.abc
from pptx import Presentation
from pptx.util import Inches
import os
from simple_image_download import simple_image_download as simp
from PIL import Image
import random
from src.utils import *

DATA_FOLDER = r"data"
REPONSE_FILE = os.path.join(DATA_FOLDER, "output2.txt")
TEMPLATE_PPTX = os.path.join(DATA_FOLDER, "samples_template_2.pptx")
OUTPUT_PPTX = os.path.join(DATA_FOLDER, "output_slide_json2.pptx")

FONT_FOLDER = r"fonts"
CHOSEN_FONT = os.path.join(FONT_FOLDER, "Calibri Regular.ttf")

content_json = create_content_json(REPONSE_FILE)
if content_json:
    content_json = json.loads(content_json)
else:
    print("Cannot extract json from text. Please try again !!!")

try:
    prs = Presentation(TEMPLATE_PPTX)
    print(f"Use template from {TEMPLATE_PPTX}")
except:
    prs = Presentation()
    print(f"Cannot use template from {TEMPLATE_PPTX}. Create a blank file.")

#################################### Delete all existed slides ####################################
for i in range(len(prs.slides) - 1, -1, -1):
    rId = prs.slides._sldIdLst[i].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[i]

########################################## Create slides #########################################
key = list(content_json.keys())[0]
for item in content_json[key]:
    header, content = item["header"], item["content"]
    query = header.replace(" ", "_")
    image = None
    if query != "Introduction":
        downloader = simp.Downloader()
        downloader.download(query, 10)
        image_names = os.listdir(os.path.join("simple_images", query))
        while True:
            path_to_image = os.path.join(
                "simple_images",
                query,
                image_names[random.randint(0, len(image_names) - 1)],
            )
            image = Image.open(path_to_image)
            if image.size != (80, 36):
                break

    slide_layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = header

    for i, place_holder in enumerate(slide.placeholders):
        if i < 2:
            continue
        sp = place_holder.element
        sp.getparent().remove(sp)

    if image:
        picture = slide.shapes.add_picture(path_to_image, Inches(1), Inches(1))
        # TODO: change picture size

    left, top, width, height = (Inches(1.5), Inches(1.5), Inches(6), Inches(6))
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = content
    text_frame.fit_text(font_file=CHOSEN_FONT, max_size=22)
    text_frame.word_wrap = True

prs.save(OUTPUT_PPTX)
