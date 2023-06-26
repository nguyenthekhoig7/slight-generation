import io
import os
import gc
import urllib
import random
import numpy as np
from PIL import Image
from typing import Optional
from fastapi import FastAPI
from pptx.util import Inches
from pptx import Presentation
from pydantic import BaseModel
from fastapi.responses import FileResponse
from bing_image_downloader import downloader
# from bing_image_urls import bing_image_urls
from fastapi.middleware.cors import CORSMiddleware

from src.utils import *
from src.text_gen import *
from api_key import *

import shutil

DATA_FOLDER = r"data"
FONT_FOLDER = r"fonts"
TEMPLATE_FOLDER = r"pptx_templates"
template_names = os.listdir(TEMPLATE_FOLDER)
os.makedirs(DATA_FOLDER, exist_ok=True)

IMAGE_FOLDER = os.path.join(r"images")
TEMPLATE_PPTX = os.path.join(DATA_FOLDER, "template.pptx")
CHOSEN_FONT = os.path.join(FONT_FOLDER, "Calibri Regular.ttf")

app = FastAPI()

# middlewares
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# import nest_asyncio
# from pyngrok import ngrok
# import uvicorn

# ngrok.set_auth_token(NGROK_TOKEN)

# # specify a port
# port = 8000
# ngrok_tunnel = ngrok.connect(port)

# # where we can visit our fastAPI app
# print('Public URL:', ngrok_tunnel.public_url)


# nest_asyncio.apply()


class Input(BaseModel):
    topic: str
    mode: int = 0
    n_slides: Optional[int] = 10
    n_words_per_slide: Optional[int] = 70
    api_token: str = None


@app.post("/generate/")
async def generate(inp_params: Input):
    if inp_params.api_token is None:
        inp_params.api_token = POE_API_KEY
    if inp_params.mode == 0:
        text_query = create_query_from_text(
            inp_params.topic, type_of_text="topic", n_slides=inp_params.n_slides, n_words_per_slide=inp_params.n_words_per_slide
        )
    else:
        text_query = create_query_from_text(
            inp_params.topic, type_of_text="doc", n_slides=inp_params.n_slides, n_words_per_slide=inp_params.n_words_per_slide
        )

    response = query_from_API(query=text_query, token=inp_params.api_token)
    content_json = create_content_json(response)

    if content_json is None:
        print("Cannot extract json from text. Cannot create a presentation. Stopped.")
        exit()

    TEMPLATE_PPTX = os.path.join(TEMPLATE_FOLDER, random.choice(template_names))
    try:
        default_16_9_slide_size = (Inches(5.625), Inches(10))
        prs = Presentation(TEMPLATE_PPTX)
        if not (
            prs.slide_height == default_16_9_slide_size[0] * 914400 and prs.slide_width == default_16_9_slide_size[1] * 914400
        ):
            print(f"Use template from {TEMPLATE_PPTX}")
        else:
            prs = Presentation()
            prs.slide_height, prs.slide_width = default_16_9_slide_size
            print("Template is not of 16:9 ratio, creating from blank template.")

    except:
        prs = Presentation()
        prs.slide_height, prs.slide_width = default_16_9_slide_size
        print(f"Cannot use template from {TEMPLATE_PPTX}. Creating a blank file.")

    layout_id = get_layout_id(prs)

    for i in range(len(prs.slides) - 1, -1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]

    # get_title_query = create_query_get_title(content_json)
    # title = query_from_API(query=get_title_query, token=POE_API_KEY)
    title = inp_params.topic
    slide_title = prs.slides.add_slide(prs.slide_layouts[0])
    title_box = slide_title.shapes.title
    title_box.text = title.upper()
    for i, place_holder in enumerate(slide_title.placeholders):
        if i < 1:
            continue
        sp = place_holder.element
        sp.getparent().remove(sp)

    key = list(content_json.keys())[0]

    if not os.path.isdir(IMAGE_FOLDER):
        os.makedirs(IMAGE_FOLDER)

    img_slot = {"left": Inches(6), "top": Inches(1.6), "width": Inches(3.6), "height": Inches(3.6)}
    img_slot_ratio = img_slot["width"] / img_slot["height"]

    for item in content_json[key]:
        header, content = process_header(item["header"]), item["content"]
        image_query = f"{header}_{inp_params.topic}".replace(" ", "_")
        image = None
        if "Introduction" not in image_query:
            try:
                downloader.download(
                    image_query,
                    limit=1,
                    output_dir=IMAGE_FOLDER,
                    force_replace=False,
                    adult_filter_off=False,
                    filter="photo",
                    timeout=10,
                    verbose=False,
                )

                image_folder_path = os.path.join(IMAGE_FOLDER, image_query)
                image_path = os.path.join(image_folder_path, os.listdir(image_folder_path)[0])
                image = Image.open(image_path)
                print(image_path)
                
                # image_urls = bing_image_urls(image_query, limit=10)
                # random.shuffle(image_urls)
                # for url in image_urls:
                #     try:
                #         image_content = urllib.request.urlopen(url, timeout=3)
                #         image_pil = Image.open(image_content)
                #         image_out = io.BytesIO()
                #         image_pil.save(image_out, format="PNG")
                #         break
                #     except:
                #         print("Cannot download image from url: {}".format(url))
                #         continue

                print('Downloaded image for slide "{}"'.format(header))
            except Exception as e:
                print(e)

        slide_layout = prs.slide_layouts[layout_id]
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = header

        for i, place_holder in enumerate(slide.placeholders):
            if i < 1:
                continue
            sp = place_holder.element
            sp.getparent().remove(sp)

        if image:
            w, h = image.size
            try:
                if w / h > img_slot_ratio:  # image longer than slot -> resize image to fit slot_width
                    slide.shapes.add_picture(image_path, img_slot["left"], img_slot["top"], width=img_slot["width"])
                else:
                    slide.shapes.add_picture(image_path, img_slot["left"], img_slot["top"], height=img_slot["height"])
                shutil.rmtree(image_folder_path)
                del image
                
            except Exception as e:
                print("Cannot add picture. ", e)

        left, top, width, height = Inches(1), Inches(1.5), Inches(5), Inches(5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = content
        text_frame.fit_text(font_file=CHOSEN_FONT, max_size=22)
        text_frame.word_wrap = True


    output_pptx_path = os.path.join("data", inp_params.topic.replace(" ", "_") + ".pptx")
    output_pptx_path = change_name_if_duplicated(output_pptx_path)
    prs.save(output_pptx_path)

    # output_folder = os.path.splitext(output_pptx_path)[-2]
    # os.makedirs(output_folder)
    # convert_pptx_to_svg(pptx_file=output_pptx_path, output_folder=output_folder)
    print(f"Presentation saved to {output_pptx_path}")

    return FileResponse(
        output_pptx_path,
        filename=output_pptx_path.split("/")[-1],
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


# # finally run the app
# uvicorn.run(app, port=port)
