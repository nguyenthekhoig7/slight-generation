import streamlit as st
import os
import json
import random
from PIL import Image
import collections.abc
from pptx.util import Inches
from pptx import Presentation
from api_key import POE_API_KEY
from io import StringIO

from src.utils import *
from src.image_download import Downloader
from src.text_gen import *
import re

DATA_FOLDER = r"data"
FONT_FOLDER = r"fonts"
TEMPLATE_PPTX = os.path.join(DATA_FOLDER, "template.pptx")
CHOSEN_FONT = os.path.join(FONT_FOLDER, "Calibri Regular.ttf")


def _create_content_from_json(response):
    match = re.search(r"{(.*?)]\n}", response, re.DOTALL)
    if match:  # response has json inside
        content_json = match.group(0)
        try:
            content_json = json.loads(content_json)
            return content_json
        except:
            return None


def _create_content_from_python_code(response):
    match_pycode = re.search(r"```python(.*?)\)\n```", response, re.DOTALL)
    if not match_pycode:
        return None
    py_code = match_pycode.group(0)

    match_content = re.search(r"content =(.*?)]", py_code, re.DOTALL)
    if not match_content:
        return None
    content_str = match_content.group(0)
    contents = content_str.split(r"=")[1]
    contents_list = re.findall('"(.*)",', contents)

    match_header = re.search(r"header =(.*?)]", response, re.DOTALL)
    if not match_header:
        return None
    headers_str = match_header.group(0)
    headers_list = re.findall('"Slide \d+: (.+)"', headers_str)

    slides_json = {}
    slides_json["slides"] = []

    for header, content in zip(headers_list, contents_list):
        pair = {}
        pair["header"] = header
        pair["content"] = content
        slides_json["slides"].append(pair)

    return slides_json


downloader = Downloader()

tab_1, tab_2 = st.tabs(["Upload document", "Input topic"])
with tab_1:
    docu_file = st.file_uploader("Choose a document", type=["docx", "pdf", "doc", "txt"])
    n_slides, n_words_per_slide = 10, 55
    # n_slides = st.number_input("Number of slide")
    # n_words_per_slide = st.number_input("Number of words per slide")
    if docu_file:
        stringio = StringIO(docu_file.getvalue().decode("utf-8"))
        string_data = stringio.read()
        st.write(string_data)
        query = """"{
        "input_text": "[[QUERY]]",
        "output_format": "json",
        "json_structure": {
            "slides":"{{presentation_slides}}"
        }
        }"""

        topic_query = (
            f"Generate a {n_slides} slide presentation from the document provided. Produce {n_words_per_slide-5} to {n_words_per_slide+5} words per slide. "
            + ". Each slide should have a  {{header}}, {{content}}. The first slide should only contain the short title. The final slide should be some discussion questions, seperated by a newline character. Return as JSON, only JSON, not the code to generate JSON."
            + " Here is the document: \n"
            + string_data
        )
        text_query = query.replace("[[QUERY]]", topic_query)
        get_title_query = (
            "Summarize the following document into a short title. Your response should contain only a short title of less than 10 words, nothing other than that, a less-than-10-word title. Here is the document: \n"
            + string_data
        )
        st.write("get_title_query" + get_title_query)
        # topic = query_from_API(query=get_title_query, token=POE_API_KEY)
        response = query_from_API(text_query, token=POE_API_KEY, bot_name="chinchilla")
        st.write("Response: " + response)
        content_json = _create_content_from_json(response)
        if content_json is None:
            content_json = _create_content_from_python_code(response)
        st.write(content_json)


with tab_2:
    topic = st.text_input("What do you want to make a presentation about?")
    st.write(topic)


# if mode[input_mode] == "document":
#     docu_file = input("Enter the document(docx/pdf) path:\n >>> ")
#     while True:
#         if ".doc" in docu_file or '.pdf' in docu_file:
#             break
#         print("Only accept .doc or .docx files. Please try again.")
#         docu_file = input("Enter the document(docx/pdf) path:\n >>> ")
#     text_query, output_txt_path = create_query_read_document(docu_file=docu_file)
#     # output_txt_path = os.path.join("data", re.sub(r"\.[a-zA-Z0-9]+$", ".txt", docu_file))  # .docx and .doc --> .txt

#     with open(output_txt_path, "r") as f:
#         docu = f.read()
#     get_title_query = (
#         "Summarize the following document into a short title. Your response should contain only a short title of less than 10 words, nothing other than that, a less-than-10-word title. Here is the document: \n"
#         + docu
#     )
#     topic = query_from_API(query=get_title_query, token=POE_API_KEY)

# else:  # mode topic
#     topic = input("What do you want to make a presentation about? \n >>> ")
#     text_query = create_query(topic, n_slides=10, n_words_per_slide=70)
#     output_txt_path = os.path.join("data", topic.replace(" ", "_") + ".txt")

# success = query_API__save_to_file(query=text_query, token=POE_API_KEY, output_path=output_txt_path)

# if success:
#     print(f"Successfully generate content about {topic}.")
# else:
#     print("Cannot query from API. Please try again")

# content_json = create_content_json(output_txt_path)
# if content_json is None:
#     print("Cannot extract json from text. Please try again !!!")

# try:
#     prs = Presentation(TEMPLATE_PPTX)
#     print(f"Use template from {TEMPLATE_PPTX}")
# except:
#     prs = Presentation()
#     print(f"Cannot use template from {TEMPLATE_PPTX}. Creating a blank file.")

# for i in range(len(prs.slides) - 1, -1, -1):
#     rId = prs.slides._sldIdLst[i].rId
#     prs.part.drop_rel(rId)
#     del prs.slides._sldIdLst[i]
# try:
#     key = list(content_json.keys())[0]
# except:
#     key = None
#     print('`content_json` is empty.')
# for item in content_json[key]:
#     header, content = process_header(item["header"]), item["content"]
#     image_query = (header + topic).replace(" ", "_")
#     image = None
#     if "Introduction" not in image_query:
#         try:
#             downloader.download(image_query, limit=20, timer=50)
#             image_names = os.listdir(os.path.join("simple_images", image_query))
#             while True:
#                 path_to_image = os.path.join(
#                     "simple_images",
#                     image_query,
#                     image_names[random.randint(0, len(image_names) - 1)],
#                 )
#                 image = Image.open(path_to_image)
#                 if image.size != (80, 36) and check_valid_image(path_to_image):
#                     break
#             downloader.flush_cache()
#         except:
#             print("Cannot download image")

#     slide_layout = prs.slide_layouts[3]
#     slide = prs.slides.add_slide(slide_layout)

#     title = slide.shapes.title
#     title.text = header

#     for i, place_holder in enumerate(slide.placeholders):
#         if i < 1:
#             continue
#         sp = place_holder.element
#         sp.getparent().remove(sp)

#     if image:
#         w, h = image.size
#         if w > h:
#             picture = slide.shapes.add_picture(path_to_image, Inches(6), Inches(2.5), width=Inches(3.8))
#         else:
#             picture = slide.shapes.add_picture(path_to_image, Inches(6), Inches(2.5), height=Inches(5))

#     left, top, width, height = Inches(1), Inches(2.5), Inches(5), Inches(5)
#     textbox = slide.shapes.add_textbox(left, top, width, height)
#     text_frame = textbox.text_frame
#     text_frame.text = content
#     text_frame.fit_text(font_file=CHOSEN_FONT, max_size=22)
#     text_frame.word_wrap = True

# output_pptx_path = output_txt_path.replace("txt", "pptx")
# prs.save(output_pptx_path)
