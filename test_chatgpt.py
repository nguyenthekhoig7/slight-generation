import openai, json
from pptx import Presentation
openai.api_key = "YOUR OPENAI API-KEY HERE!"

presentation_title = input("What do you want to make a presentation about?")

query_json = """"{
    "input_text": "[[QUERY]]",
    "output_format": "json",
    "json_structure": {
        "slides":"{{presentation_slides}}"
       }
    }"""
       
question = "Generate a 10 slide presentation for the topic. Produce 50 to 60 words per slide. " + presentation_title + ".Each slide should have a  {{header}}, {{content}}. The final slide should be a list of discussion questions. Return as JSON."

prompt = query_json.replace("[[QUERY]]",question)
print(prompt)
completion = openai.ChatCompletion.create(model = "gpt-3.5-turbo", messages =[{"role":"user","content":prompt}])
response = completion.choices[0].message.content

print(response)

r = json.loads(response)

slide_data = r["slides"]

prs = Presentation()

for slide in slide_data:
    slide_layout = prs.slide_layouts[1]
    new_slide = prs.slides.add_slide(slide_layout)
    
    if slide['header']:
        title = new_slide.shapes.title
        title.text = slide['header']
    
    if slide['content']:
       shapes = new_slide.shapes
       body_shape = shapes.placeholders[1]
       tf = body_shape.text_frame
       tf.text = slide['content']
       tf.fit_text(font_family="Calibri", max_size=18, bold=True)
       
prs.save("output.pptx")